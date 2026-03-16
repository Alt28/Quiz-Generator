"""
CLARKO - StudyBuddy Quiz Generator  (Flask Web App)
Upload PDF / DOCX / PPTX / TXT  ->  choose quiz type  ->  get scored
"""

import os
import re
import math
import random
import tempfile
from collections import Counter

from flask import Flask, request, jsonify, render_template

# ── NLP ──
import spacy

# ── file parsers ──
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

# Load spaCy model once at startup
nlp = spacy.load("en_core_web_sm")
nlp.max_length = 1_500_000  # allow large documents

app = Flask(
    __name__,
    template_folder=os.path.join(os.path.dirname(__file__), "templates"),
    static_folder=os.path.join(os.path.dirname(__file__), "static"),
)
app.config["MAX_CONTENT_LENGTH"] = 20 * 1024 * 1024  # 20 MB

ALLOWED_EXTENSIONS = {"pdf", "docx", "pptx", "txt"}

# ──────────────────────────────────────────────
#  STOPWORDS  (supplement spaCy's built-in list)
# ──────────────────────────────────────────────
EXTRA_STOPS = set("""
also used using many much however may well even still just like way one two three
make made use new include includes including example examples called known type types
can based also step note please next first second third finally click select go
configure according""".split())


# ══════════════════════════════════════════════
#  FILE TEXT EXTRACTION
# ══════════════════════════════════════════════

def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_text_from_pdf(path: str) -> str:
    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                parts.append(t)
    return "\n".join(parts)


def extract_text_from_docx(path: str) -> str:
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


def extract_text(path: str, ext: str) -> str:
    if ext == "pdf":
        return extract_text_from_pdf(path)
    if ext == "docx":
        return extract_text_from_docx(path)
    if ext == "pptx":
        return extract_text_from_pptx(path)
    if ext == "txt":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    return ""


# ══════════════════════════════════════════════
#  TEXT PREPROCESSING  (spaCy-powered)
# ══════════════════════════════════════════════

def clean_text(text: str) -> str:
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\w\s.,;:!?\'\-/()"]', '', text)
    return text.strip()


def _is_stop(token) -> bool:
    """Check if a spaCy token is a stopword (built-in + custom)."""
    return token.is_stop or token.lemma_.lower() in EXTRA_STOPS


def _make_doc(text: str):
    """Run spaCy pipeline on text, return Doc object."""
    return nlp(text)


def split_sentences(text: str, doc=None) -> list[str]:
    """Split text into sentences using spaCy's sentence boundary detection."""
    if doc is None:
        doc = _make_doc(text)
    result = []
    for sent in doc.sents:
        s = sent.text.strip()
        word_count = len(s.split())
        if len(s) > 15 and 4 <= word_count <= 35:
            result.append(s)
    return result


def tokenize(text: str) -> list[str]:
    """Lemmatize content words using spaCy."""
    doc = _make_doc(text)
    return [
        tok.lemma_.lower()
        for tok in doc
        if tok.is_alpha and len(tok.text) > 2 and not _is_stop(tok)
    ]


def tokenize_doc(doc) -> list[str]:
    """Lemmatize content words from an existing spaCy Doc (no re-parse)."""
    return [
        tok.lemma_.lower()
        for tok in doc
        if tok.is_alpha and len(tok.text) > 2 and not _is_stop(tok)
    ]


def tokenize_fast(text: str) -> set:
    """Quick lowercase word set for scoring (no spaCy, just regex split)."""
    return set(w.lower() for w in re.findall(r'[A-Za-z]{3,}', text))


# ══════════════════════════════════════════════
#  KEY CONCEPT EXTRACTION  (spaCy NLP + TF-IDF)
# ══════════════════════════════════════════════

def _idf(word: str, sentence_tokens: list[list[str]]) -> float:
    """Inverse document frequency across sentences."""
    n = len(sentence_tokens) or 1
    df = sum(1 for sent_toks in sentence_tokens if word in sent_toks)
    return math.log((n + 1) / (df + 1)) + 1


def extract_keywords(text: str, top_n: int = 25, doc=None) -> list[str]:
    """Score words using spaCy POS tags, NER, lemma frequency, and TF-IDF.
    Returns original surface forms (best casing) for better blank matching."""
    if doc is None:
        doc = _make_doc(text)

    # Collect content-word lemmas (NOUN, PROPN, ADJ)
    content_pos = {'NOUN', 'PROPN', 'ADJ'}
    lemmas = []
    # Track best surface form per lemma (prefer capitalized/original)
    surface_forms = {}
    propn_lemmas = set()
    for tok in doc:
        if (tok.pos_ in content_pos and tok.is_alpha
                and len(tok.text) > 2 and not _is_stop(tok)):
            lemma = tok.lemma_.lower()
            lemmas.append(lemma)
            # Keep the most informative surface form
            if lemma not in surface_forms or tok.pos_ == 'PROPN':
                surface_forms[lemma] = tok.text
            if tok.pos_ == 'PROPN':
                propn_lemmas.add(lemma)

    freq = Counter(lemmas)
    total = len(lemmas) or 1

    # Build per-sentence token sets for IDF
    sent_tokens = []
    for sent in doc.sents:
        toks = set()
        for tok in sent:
            if tok.pos_ in content_pos and tok.is_alpha and not _is_stop(tok):
                toks.add(tok.lemma_.lower())
        sent_tokens.append(toks)

    # Named-entity bonus
    ent_lemmas = set()
    for ent in doc.ents:
        if ent.label_ in ('ORG', 'PRODUCT', 'EVENT', 'WORK_OF_ART', 'LAW',
                          'GPE', 'LOC', 'PERSON', 'NORP', 'FAC'):
            for tok in ent:
                if tok.is_alpha and not _is_stop(tok):
                    ent_lemmas.add(tok.lemma_.lower())

    # Words near copula definitions get a boost
    def_lemmas = set()
    for tok in doc:
        if tok.dep_ == 'nsubj' and tok.head.lemma_ == 'be':
            if tok.is_alpha and tok.pos_ in content_pos:
                def_lemmas.add(tok.lemma_.lower())

    scored = {}
    for word, count in freq.items():
        tf = count / total
        idf = _idf(word, sent_tokens)
        ent_bonus = 0.4 if word in ent_lemmas else 0
        def_bonus = 0.5 if word in def_lemmas else 0
        pos_bonus = 0.2 if word in propn_lemmas else 0
        scored[word] = tf * idf * 100 + ent_bonus + def_bonus + pos_bonus

    ranked = sorted(scored, key=scored.get, reverse=True)
    # Return surface forms for better text matching
    return [surface_forms.get(w, w) for w in ranked[:top_n]]


def extract_noun_phrases(text: str, doc=None) -> list[str]:
    """Extract multi-word noun phrases using spaCy noun chunks + NER."""
    if doc is None:
        doc = _make_doc(text)
    seen = set()
    phrases = []

    # Named entities first (highest quality)
    for ent in doc.ents:
        phrase = ent.text.strip()
        if (len(phrase.split()) >= 1 and phrase.lower() not in seen
                and len(phrase) > 2 and phrase.lower() not in EXTRA_STOPS):
            seen.add(phrase.lower())
            phrases.append(phrase)

    # Noun chunks (multi-word only, skip determiners)
    for chunk in doc.noun_chunks:
        # Remove leading determiners/pronouns
        tokens = [tok for tok in chunk if tok.pos_ not in ('DET', 'PRON')]
        if not tokens:
            continue
        phrase = ' '.join(tok.text for tok in tokens)
        if (len(tokens) >= 2 and phrase.lower() not in seen
                and phrase.lower() not in EXTRA_STOPS):
            seen.add(phrase.lower())
            phrases.append(phrase)

    return phrases[:30]


# ══════════════════════════════════════════════
#  DEFINITION & RELATIONSHIP DETECTION  (spaCy dep parse + regex fallback)
# ══════════════════════════════════════════════

# Regex fallback patterns (kept for texts where dep parsing misses things)
DEF_PATTERNS = [
    re.compile(
        r'^([A-Z][A-Za-z\s]{1,40}?)\s+(?:is|are)\s+(?:a|an|the)?\s*(.{10,})',
        re.IGNORECASE
    ),
    re.compile(
        r'^([A-Z][A-Za-z\s]{1,40}?)\s+refers?\s+to\s+(.{10,})',
        re.IGNORECASE
    ),
    re.compile(
        r'^([A-Z][A-Za-z\s]{1,40}?)\s+is\s+(?:defined|known|described)\s+as\s+(.{10,})',
        re.IGNORECASE
    ),
    re.compile(
        r'^([A-Z][A-Za-z\s]{1,40}?)\s+means?\s+(.{10,})',
        re.IGNORECASE
    ),
    re.compile(
        r'^([A-Z][A-Za-z\s]{1,40}?)\s+(?:consists?\s+of|involves?|includes?)\s+(.{10,})',
        re.IGNORECASE
    ),
]

BAD_STARTS = {'even', 'if', 'when', 'while', 'although', 'because', 'since',
              'after', 'before', 'unless', 'until', 'once', 'make', 'step',
              'note', 'please', 'however', 'therefore', 'thus', 'also',
              'then', 'next', 'first', 'second', 'third', 'finally',
              'it', 'they', 'he', 'she', 'we', 'you', 'this', 'that',
              'these', 'those', 'there', 'here', 'its', 'their'}

_LEADING_ARTICLES = re.compile(r'^(?:a|an|the)\s+', re.IGNORECASE)


def _normalize_term(term: str) -> str:
    """Strip leading articles and extra whitespace from a term."""
    term = _LEADING_ARTICLES.sub('', term).strip()
    return term


def _extract_defs_spacy(doc) -> list[dict]:
    """Use spaCy dependency parsing to find 'X is Y' definition patterns."""
    defs = []
    seen = set()
    for sent in doc.sents:
        # Look for copula pattern: nsubj ← be → attr/acomp
        root = sent.root
        if root.lemma_ != 'be':
            continue
        # Find the subject
        subj = None
        for child in root.children:
            if child.dep_ in ('nsubj', 'nsubjpass'):
                subj = child
                break
        if not subj:
            continue
        # Find the attribute (what the subject "is")
        attr = None
        for child in root.children:
            if child.dep_ in ('attr', 'acomp', 'oprd'):
                attr = child
                break
        if not attr:
            continue

        # Build term from subject + its subtree (noun chunk)
        subj_span = doc[subj.left_edge.i: subj.right_edge.i + 1]
        term_tokens = [t for t in subj_span
                       if t.pos_ not in ('DET', 'PRON', 'ADP', 'PUNCT')]
        if not term_tokens:
            continue
        # Skip if subject is a pronoun (e.g. "It is ...")
        if subj.pos_ == 'PRON':
            continue
        term = ' '.join(t.text for t in term_tokens)
        term = _normalize_term(term)

        # Skip bad terms
        words = term.split()
        if len(words) > 4 or len(words) == 0:
            continue
        if words[0].lower() in BAD_STARTS:
            continue
        if re.search(r'\d|step|configure|make sure|click|select|go to', term, re.I):
            continue

        # Build definition from attr onward
        attr_span = doc[attr.left_edge.i: sent.end]
        definition = attr_span.text.strip().rstrip('.')
        if len(definition.split()) < 3:
            continue

        if term.lower() not in seen:
            seen.add(term.lower())
            defs.append({
                'term': term,
                'definition': definition,
                'sentence': sent.text.strip(),
            })
    return defs


def extract_definitions(sentences: list[str], doc=None) -> list[dict]:
    """Find term→definition pairs using spaCy dep parsing + regex fallback."""
    defs = []
    seen_terms = set()

    # Phase 1: spaCy dependency parsing (high quality)
    if doc is not None:
        spacy_defs = _extract_defs_spacy(doc)
        for d in spacy_defs:
            if d['term'].lower() not in seen_terms:
                seen_terms.add(d['term'].lower())
                defs.append(d)

    # Phase 2: regex fallback for patterns spaCy might miss
    for sent in sentences:
        for pat in DEF_PATTERNS:
            m = pat.match(sent)
            if m:
                term = _normalize_term(m.group(1).strip().rstrip(','))
                definition = m.group(2).strip().rstrip('.')
                words = term.split()
                if len(words) > 4:
                    break
                if words[0].lower() in BAD_STARTS:
                    break
                if re.search(r'\d|step|configure|make sure|click|select|go to',
                             term, re.IGNORECASE):
                    break
                if term.lower() not in seen_terms:
                    seen_terms.add(term.lower())
                    defs.append({
                        'term': term,
                        'definition': definition,
                        'sentence': sent,
                    })
                break
    return defs


def detect_lists(text: str) -> list[dict]:
    """Detect numbered/bulleted lists in the text for enumeration questions."""
    lists = []

    # Numbered lists: "1. item", "2. item", etc.
    numbered = re.findall(r'(?:^|\n)\s*\d+[.)]\s*(.+)', text)
    if len(numbered) >= 3:
        lists.append({
            'question': 'Enumerate the items from the list in the notes.',
            'items': [item.strip().rstrip('.') for item in numbered],
        })

    # Dash/bullet lists: "- item", "• item", "* item"
    bulleted = re.findall(r'(?:^|\n)\s*[-•*]\s+(.+)', text)
    if len(bulleted) >= 3:
        lists.append({
            'question': 'Enumerate the items from the bulleted list.',
            'items': [item.strip().rstrip('.') for item in bulleted],
        })

    # Semicolon lists: "A; B; C"
    for sent in text.split('.'):
        parts = [p.strip() for p in sent.split(';') if len(p.strip()) > 2]
        if len(parts) >= 3:
            lists.append({
                'question': f'Enumerate the items: "{sent.strip()[:70]}..."',
                'items': parts,
            })
            break

    return lists


def extract_key_sentences(text: str, keywords: list[str], top_n: int = 15,
                          doc=None) -> list[str]:
    """Rank sentences by keyword density, entity richness, and information quality."""
    sentences = split_sentences(text, doc=doc)
    kw_set = set(k.lower() for k in keywords)

    # Pre-compute sentence-level entity counts from the spaCy doc
    sent_ent_count = {}
    if doc is not None:
        for sent in doc.sents:
            s = sent.text.strip()
            ent_count = sum(1 for ent in sent.ents if ent.label_ not in ('CARDINAL', 'ORDINAL'))
            sent_ent_count[s] = ent_count

    scored = []
    for sent in sentences:
        # Use fast regex tokenizer instead of full spaCy reparse
        words = tokenize_fast(sent)
        kw_overlap = len(words & kw_set)
        # Bonus for definition-like sentences
        def_bonus = 2 if any(pat.match(sent) for pat in DEF_PATTERNS) else 0
        # Entity density bonus
        ent_bonus = sent_ent_count.get(sent, 0) * 0.5
        # Penalty for very short or very long
        length = len(sent.split())
        length_score = 1.0 if 8 <= length <= 40 else 0.5
        score = (kw_overlap + def_bonus + ent_bonus) * length_score
        scored.append((score, sent))

    scored.sort(key=lambda x: x[0], reverse=True)
    return [sent for _, sent in scored[:top_n]]


# ══════════════════════════════════════════════
#  QUESTION GENERATORS
# ══════════════════════════════════════════════

def generate_fill_in_the_blank(sentences, keywords, definitions, max_q=10):
    """Create fill-in-the-blank questions prioritizing definition sentences."""
    questions = []
    kw_lower = {k.lower() for k in keywords}
    used_pairs = set()  # track (sentence, keyword) pairs to allow reuse with different blanks

    # Priority 1: definition sentences — blank the term
    for d in definitions:
        sent = d['sentence']
        term = d['term']
        pair_key = (sent, term.lower())
        if pair_key in used_pairs:
            continue
        blanked = re.sub(
            r'\b' + re.escape(term) + r'\b',
            '________',
            sent,
            count=1,
            flags=re.IGNORECASE,
        )
        if blanked != sent:
            used_pairs.add(pair_key)
            questions.append({'question': blanked, 'answer': term})
        if len(questions) >= max_q:
            return questions

    # Priority 2: keyword-rich sentences (can blank multiple keywords per sentence)
    kw_rank = {k.lower(): i for i, k in enumerate(keywords)}
    for sent in sentences:
        if len(questions) >= max_q:
            break
        if len(sent.split()) > 30:
            continue
        words_in_sent = re.findall(r'\b[A-Za-z]{3,}\b', sent)
        targets = [w for w in words_in_sent if w.lower() in kw_lower]
        # Sort by keyword rank (best first)
        targets.sort(key=lambda w: kw_rank.get(w.lower(), 999))
        for target in targets:
            if len(questions) >= max_q:
                break
            pair_key = (sent, target.lower())
            if pair_key in used_pairs:
                continue
            blanked = re.sub(
                r'\b' + re.escape(target) + r'\b',
                '________',
                sent,
                count=1,
                flags=re.IGNORECASE,
            )
            if blanked != sent:
                used_pairs.add(pair_key)
                questions.append({'question': blanked, 'answer': target})
    return questions


def _pick_distractors(answer: str, keywords: list[str], definitions: list[dict],
                      count: int = 3) -> list[str]:
    """Pick plausible distractors — prefer semantically related terms."""
    answer_lower = answer.lower().strip()
    answer_word_count = len(answer.split())

    # Collect candidate terms
    candidates = []
    for k in keywords:
        k = k.strip()
        if k.lower() == answer_lower:
            continue
        if len(k.split()) > 3:
            continue
        if any(w in k.lower() for w in ['make sure', 'that the', 'in order', 'such as',
                                          'as well', 'due to', 'based on', 'according']):
            continue
        candidates.append(k)
    for d in definitions:
        t = d['term'].strip()
        if t.lower() == answer_lower or t.lower() in [c.lower() for c in candidates]:
            continue
        if len(t.split()) > 3:
            continue
        candidates.append(t)

    if not candidates:
        return []

    # Score using spaCy word vectors (semantic similarity) + surface heuristics
    answer_doc = nlp(answer)
    answer_len = len(answer)
    scored = []
    for c in candidates:
        c_doc = nlp(c)
        # Semantic similarity (0-1, higher = more similar = better distractor)
        sim = answer_doc.similarity(c_doc) if answer_doc.has_vector and c_doc.has_vector else 0
        len_diff = abs(len(c) - answer_len)
        wc_diff = abs(len(c.split()) - answer_word_count)
        case_match = 1 if (c[0].isupper() == answer[0].isupper()) else 0
        # Similarity is most important, then surface similarity
        score = sim * 10 + case_match * 2 - len_diff * 0.05 - wc_diff * 1
        scored.append((score, c))

    scored.sort(key=lambda x: x[0], reverse=True)
    result = [c for _, c in scored[:count * 2]]
    random.shuffle(result)
    return result[:count]


def generate_mcq(sentences, keywords, definitions, max_q=10):
    """Create MCQs with short, clean term choices only."""
    questions = []
    used_pairs = set()  # (sentence, keyword) pairs

    # Style 1: Definition-based — show definition as question, terms as choices
    if len(definitions) >= 4:
        for d in definitions:
            if len(questions) >= max_q:
                break
            term = d['term']
            definition = d['definition']
            if len(definition.split()) < 4:
                continue
            # Truncate definition for readable question
            defn_display = definition if len(definition) <= 100 else definition[:97] + '...'
            # Other definition terms as distractors (short terms only)
            other_terms = [
                dd['term'] for dd in definitions
                if dd['term'].lower() != term.lower()
                and len(dd['term'].split()) <= 4
            ]
            random.shuffle(other_terms)
            dists = other_terms[:3]
            if len(dists) < 3:
                continue
            choices = dists + [term]
            random.shuffle(choices)
            used_pairs.add((d['sentence'], term.lower()))
            questions.append({
                'question': defn_display,
                'choices': choices,
                'answer': term,
            })

    # Style 2: Sentence-blank MCQ — answer is always a single keyword
    kw_lower = {k.lower() for k in keywords}
    kw_index = {k.lower(): i for i, k in enumerate(keywords)}
    for sent in sentences:
        if len(questions) >= max_q:
            break
        # Skip sentences that are too long for a readable question
        if len(sent.split()) > 25:
            continue
        words_in_sent = re.findall(r'\b[A-Za-z]{3,}\b', sent)
        targets = [w for w in words_in_sent if w.lower() in kw_lower]
        if not targets:
            continue
        targets.sort(key=lambda w: kw_index.get(w.lower(), 999))
        for target in targets:
            if len(questions) >= max_q:
                break
            pair_key = (sent, target.lower())
            if pair_key in used_pairs:
                continue
            blanked = re.sub(
                r'\b' + re.escape(target) + r'\b',
                '________',
                sent,
                count=1,
                flags=re.IGNORECASE,
            )
            if blanked == sent:
                continue
            distractors = _pick_distractors(target, keywords, definitions)
            if len(distractors) < 3:
                continue
            # Verify all choices are clean short terms
            all_clean = all(len(c.split()) <= 3 for c in distractors + [target])
            if not all_clean:
                continue
            choices = distractors + [target]
            random.shuffle(choices)
            used_pairs.add(pair_key)
            questions.append({
                'question': blanked,
                'choices': choices,
                'answer': target,
            })

    return questions


def generate_flashcards(sentences, keywords, definitions, max_q=15):
    """Create term → definition flashcard pairs."""
    cards = []
    seen = set()

    # Priority 1: real definitions from pattern matching
    for d in definitions:
        term = d['term']
        if term.lower() in seen:
            continue
        seen.add(term.lower())
        defn = d['definition']
        if len(defn) > 150:
            defn = defn[:147] + '...'
        cards.append({'term': term, 'definition': defn})
        if len(cards) >= max_q:
            return cards

    # Priority 2: keyword → best sentence containing it
    kw_lower = {k.lower(): k for k in keywords}
    for sent in sentences:
        if len(cards) >= max_q:
            break
        for kw_l, kw_orig in kw_lower.items():
            if kw_l in sent.lower() and kw_l not in seen:
                seen.add(kw_l)
                defn = sent.strip().rstrip('.')
                if len(defn) > 150:
                    defn = defn[:147] + '...'
                cards.append({'term': kw_orig.capitalize(), 'definition': defn})
                break
    return cards


def generate_enumeration(text, sentences, keywords, definitions, max_q=6):
    """Generate enumeration questions from detected lists and grouped concepts."""
    questions = []
    used = set()

    # Priority 1: Real lists detected from the text
    detected_lists = detect_lists(text)
    for lst in detected_lists:
        if len(questions) >= max_q:
            break
        items = lst['items'][:7]  # cap at 7 items
        questions.append({
            'question': lst['question'],
            'answers': items,
            'count': len(items),
        })
        for it in items:
            used.add(it.lower())

    # Priority 2: Definition terms grouped by sentence proximity
    for sent in sentences:
        if len(questions) >= max_q:
            break
        found = []
        for d in definitions:
            if d['term'].lower() in sent.lower() and d['term'].lower() not in used:
                found.append(d['term'])
        if len(found) >= 2:
            for k in found:
                used.add(k.lower())
            short = sent[:80].rstrip() + ('...' if len(sent) > 80 else '')
            questions.append({
                'question': f'Enumerate the key terms found in: "{short}"',
                'answers': found,
                'count': len(found),
            })

    # Priority 3: Keyword groupings from sentences
    for sent in sentences:
        if len(questions) >= max_q:
            break
        found = [kw for kw in keywords if kw.lower() in sent.lower()
                 and kw.lower() not in used]
        if len(found) >= 2:
            for k in found:
                used.add(k.lower())
            short = sent[:80].rstrip() + ('...' if len(sent) > 80 else '')
            questions.append({
                'question': f'Enumerate the key terms found in: "{short}"',
                'answers': found,
                'count': len(found),
            })

    # Fallback: general key concepts
    if len(questions) == 0 and len(keywords) >= 3:
        remaining = [k for k in keywords if k.lower() not in used]
        if len(remaining) >= 3:
            chunk = remaining[:5]
            questions.append({
                'question': f'Enumerate {len(chunk)} key concepts from the given notes.',
                'answers': chunk,
                'count': len(chunk),
            })
    return questions


# ══════════════════════════════════════════════
#  MASTER QUIZ BUILDER
# ══════════════════════════════════════════════

def build_quiz(text: str, quiz_types: list[str], exam_mode: bool = False) -> dict:
    cleaned = clean_text(text)
    # Single spaCy parse — reuse across all NLP functions
    doc = _make_doc(cleaned)
    sentences = split_sentences(cleaned, doc=doc)
    kw_count = 60 if exam_mode else 25
    keywords = extract_keywords(cleaned, kw_count, doc=doc)
    noun_phrases = extract_noun_phrases(cleaned, doc=doc)
    all_keywords = list(dict.fromkeys(keywords + [p.lower() for p in noun_phrases]))
    # Filter keywords: only clean single terms (1-3 words, no fragments)
    all_keywords = [
        k for k in all_keywords
        if len(k.split()) <= 3
        and not any(frag in k.lower() for frag in [
            'make sure', 'that the', 'in order', 'such as', 'as well',
            'due to', 'based on', 'according', 'however', 'therefore',
        ])
    ]
    definitions = extract_definitions(sentences, doc=doc)
    sent_limit = len(sentences) if exam_mode else min(20, len(sentences))
    key_sentences = extract_key_sentences(cleaned, all_keywords,
                                          sent_limit, doc=doc)

    result = {
        "keywords": keywords[:15],
        "noun_phrases": noun_phrases[:10],
        "sentence_count": len(sentences),
        "definition_count": len(definitions),
        "questions": [],
        "flashcards": [],
        "exam_mode": exam_mode,
    }

    if exam_mode:
        # Exam: target 50 items total, spread across selected types
        type_count = sum(1 for t in quiz_types if t != "flashcards")
        per_type = max(50 // type_count, 15) if type_count else 50

        if "fill_in_the_blank" in quiz_types:
            for q in generate_fill_in_the_blank(key_sentences, all_keywords, definitions,
                                                 max_q=per_type):
                q["type"] = "fill_in_the_blank"
                result["questions"].append(q)

        if "multiple_choice" in quiz_types:
            for q in generate_mcq(key_sentences, all_keywords, definitions,
                                  max_q=per_type):
                q["type"] = "multiple_choice"
                result["questions"].append(q)

        if "enumeration" in quiz_types:
            for q in generate_enumeration(cleaned, key_sentences, all_keywords, definitions,
                                           max_q=per_type):
                q["type"] = "enumeration"
                result["questions"].append(q)

        if "flashcards" in quiz_types:
            result["flashcards"] = generate_flashcards(key_sentences, all_keywords, definitions,
                                                        max_q=30)
    else:
        if "fill_in_the_blank" in quiz_types:
            for q in generate_fill_in_the_blank(key_sentences, all_keywords, definitions):
                q["type"] = "fill_in_the_blank"
                result["questions"].append(q)

        if "multiple_choice" in quiz_types:
            for q in generate_mcq(key_sentences, all_keywords, definitions):
                q["type"] = "multiple_choice"
                result["questions"].append(q)

        if "enumeration" in quiz_types:
            for q in generate_enumeration(cleaned, key_sentences, all_keywords, definitions):
                q["type"] = "enumeration"
                result["questions"].append(q)

        if "flashcards" in quiz_types:
            result["flashcards"] = generate_flashcards(key_sentences, all_keywords, definitions)

    return result


# ══════════════════════════════════════════════
#  FLASK ROUTES
# ══════════════════════════════════════════════

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/generate", methods=["POST"])
def generate():
    text = ""

    # File upload
    if "file" in request.files and request.files["file"].filename:
        f = request.files["file"]
        if not allowed_file(f.filename):
            return jsonify({"error": "Unsupported file type. Use PDF, DOCX, PPTX, or TXT."}), 400
        ext = f.filename.rsplit(".", 1)[1].lower()
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix="." + ext)
        try:
            f.save(tmp.name)
            tmp.close()
            text = extract_text(tmp.name, ext)
        finally:
            os.unlink(tmp.name)

    # Pasted text fallback
    if not text.strip():
        text = request.form.get("notes", "")

    if len(text.strip()) < 30:
        return jsonify({"error": "Not enough text. Paste more notes or upload a longer document."}), 400

    quiz_types = request.form.getlist("quiz_types")
    if not quiz_types:
        quiz_types = ["multiple_choice", "fill_in_the_blank", "enumeration"]

    exam_mode = request.form.get("exam_mode") == "1"
    quiz = build_quiz(text, quiz_types, exam_mode=exam_mode)
    return jsonify(quiz)


# ══════════════════════════════════════════════
#  ENTRY POINT
# ══════════════════════════════════════════════

if __name__ == "__main__":
    app.run(debug=True, port=8080)
